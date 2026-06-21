"""scRareRefine Figure 1 — 3 columns x 2 rows PowerPoint framework.

Produces figures/scrarerefine_framework_v2.pptx with every shape independently
editable in PowerPoint (rectangles, ellipses, lines, freeform funnel/distribution).
"""
from __future__ import annotations
import math
import random
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

OUT = Path("D:/Desktop/scRareRefine/figures/scrarerefine_framework_v2.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)

# ── canvas (≈1800x1050 px @ 96dpi = 18.75x10.94 in; use widescreen 13.5x7.875) ──
SLIDE_W, SLIDE_H = 13.5, 7.875

# ── colours ────────────────────────────────────────────────────────────────────
C_MAJORITY = RGBColor(0xC7, 0xCD, 0xD4)
C_RARE     = RGBColor(0xE0, 0x80, 0x43)
C_RESCUED  = RGBColor(0x5B, 0xA8, 0x5B)
C_LATENT   = RGBColor(0x4A, 0x6F, 0xA5)
C_TAU      = RGBColor(0xC0, 0x39, 0x2B)
C_TITLE    = RGBColor(0x2C, 0x3E, 0x50)
C_BODY     = RGBColor(0x47, 0x55, 0x69)
C_MUTED    = RGBColor(0x7F, 0x8C, 0x8D)
C_FRAME    = RGBColor(0xD5, 0xD9, 0xDD)
C_PHASE    = RGBColor(0x9A, 0xA0, 0xA6)
C_AXIS     = RGBColor(0x9A, 0xA0, 0xA6)
C_ARROW    = RGBColor(0x6B, 0x72, 0x80)
C_FILL_BG  = RGBColor(0xFA, 0xFB, 0xFC)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# ── layout (inches) ────────────────────────────────────────────────────────────
M = 0.3
TITLE_Y, TITLE_H = 0.25, 0.50
PHASE_Y, PHASE_H = 0.82, 0.28
PANEL_TOP_Y = 1.22
PANEL_H = 3.00
GAP_V = 0.18
PANEL_BOT_Y = PANEL_TOP_Y + PANEL_H + GAP_V   # 4.40
COL_GAP = 0.18
COL_W = (SLIDE_W - 2 * M - 2 * COL_GAP) / 3   # ≈4.18
COL1_X = M
COL2_X = M + COL_W + COL_GAP
COL3_X = M + 2 * (COL_W + COL_GAP)


def rgb_to_xml(c: RGBColor) -> str:
    return "{:02X}{:02X}{:02X}".format(c[0], c[1], c[2])


# ── presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ── shape helpers ─────────────────────────────────────────────────────────────
def add_rect(x, y, w, h, fill=C_WHITE, line=C_FRAME, line_w=0.75, dash=False, transparent_fill=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if transparent_fill:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
        if dash:
            s.line.dash_style = 7  # MSO_LINE_DASH_STYLE.DASH = 4? value 4
    return s


def add_ellipse(cx, cy, d, fill=C_MAJORITY, line=None):
    r = d / 2.0
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(d), Inches(d))
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    return s


def add_ring(cx, cy, d, line=C_MUTED, dash=True, lw=0.75):
    r = d / 2.0
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(d), Inches(d))
    s.shadow.inherit = False
    s.fill.background()
    s.line.color.rgb = line
    s.line.width = Pt(lw)
    if dash:
        s.line.dash_style = 7
    return s


def add_line(x1, y1, x2, y2, color=C_AXIS, lw=0.75, dash=False, arrow=False):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(lw)
    if dash:
        conn.line.dash_style = 7
    if arrow:
        ln = conn.line._get_or_add_ln()
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return conn


def add_text(x, y, w, h, text, size=10, bold=False, color=C_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_x_mark(cx, cy, size=0.08, color=C_LATENT, lw=2.0):
    add_line(cx - size/2, cy - size/2, cx + size/2, cy + size/2, color=color, lw=lw)
    add_line(cx - size/2, cy + size/2, cx + size/2, cy - size/2, color=color, lw=lw)


def add_freeform(points, fill=C_MAJORITY, fill_alpha=0.55, line=C_PHASE, line_w=0.75):
    """points: list of (x,y) in inches; closes polygon."""
    builder = slide.shapes.build_freeform(Inches(points[0][0]), Inches(points[0][1]), scale=914400)
    # Convert remaining points to deltas relative to start
    x0, y0 = points[0]
    for px, py in points[1:]:
        builder.add_line_segments([(px - x0, py - y0)])
    shp = builder.convert_to_shape()
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    # alpha
    sp_fill = shp.fill._xPr.find(qn("a:solidFill")) or shp.fill._xPr.find(qn(".//a:solidFill"))
    if sp_fill is not None:
        clr = sp_fill[0]
        if "alpha" not in [child.tag.split("}")[-1] for child in clr]:
            alpha_el = etree.SubElement(clr, qn("a:alpha"))
            alpha_el.set("val", str(int(fill_alpha * 100000)))
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    return shp


# ── top title ────────────────────────────────────────────────────────────────
add_text(M, TITLE_Y, SLIDE_W - 2 * M, 0.32, "scRareRefine — framework overview", size=20, bold=True, color=C_TITLE)
add_text(M, TITLE_Y + 0.34, SLIDE_W - 2 * M, 0.20,
         "Three-phase pipeline: data representation  →  rare-cell reference and candidate selection  →  controlled rescue and final prediction",
         size=11, color=C_BODY)

# ── phase frames (dashed) + headers ─────────────────────────────────────────
for i, (cx, label_a, label_b) in enumerate([
    (COL1_X, "PHASE 1", "Data representation"),
    (COL2_X, "PHASE 2", "Rare-cell reference and candidate selection"),
    (COL3_X, "PHASE 3", "Controlled rescue and final prediction"),
]):
    add_rect(cx - 0.05, PHASE_Y - 0.02, COL_W + 0.10, (PANEL_BOT_Y + PANEL_H) - PHASE_Y + 0.04,
             transparent_fill=True, line=C_PHASE, line_w=0.75, dash=True)
    add_text(cx + 0.05, PHASE_Y, 1.0, PHASE_H, label_a, size=11, bold=True, color=C_TITLE)
    add_text(cx + 0.70, PHASE_Y, COL_W - 0.7, PHASE_H, label_b, size=10, color=C_BODY)


# ── helper: draw a panel frame and inner title ───────────────────────────────
def panel_frame(x, y, w, h, code, title, subtitle=""):
    add_rect(x, y, w, h, fill=C_WHITE, line=C_FRAME, line_w=0.75)
    add_text(x + 0.10, y + 0.05, w - 0.20, 0.22, f"{code} · {title}", size=11, bold=True, color=C_TITLE)
    if subtitle:
        add_text(x + 0.10, y + 0.27, w - 0.20, 0.18, subtitle, size=8.5, color=C_BODY)


# ===============================================================
# PANEL A  Input scRNA-seq data (matrix + label bar)
# ===============================================================
A_X, A_Y = COL1_X, PANEL_TOP_Y
panel_frame(A_X, A_Y, COL_W, PANEL_H, "A", "Input scRNA-seq data",
            "expression matrix (genes × cells)  +  sparsely labelled cells")

# matrix
M_X, M_Y, M_W, M_H = A_X + 0.55, A_Y + 0.55, COL_W - 0.75, 1.55
ROWS, COLS = 8, 24
cell_w = M_W / COLS
cell_h = M_H / ROWS
random.seed(7)
for r in range(ROWS):
    for c in range(COLS):
        alpha = 0.12 + random.random() * 0.75
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(M_X + c * cell_w), Inches(M_Y + r * cell_h),
            Inches(cell_w), Inches(cell_h),
        )
        rect.shadow.inherit = False
        rect.fill.solid()
        rect.fill.fore_color.rgb = C_LATENT
        sp_fill = rect.fill._xPr.find(qn("a:solidFill"))
        if sp_fill is not None and len(sp_fill):
            alpha_el = etree.SubElement(sp_fill[0], qn("a:alpha"))
            alpha_el.set("val", str(int(alpha * 100000)))
        rect.line.fill.background()

# axis labels
add_text(A_X + 0.10, M_Y + M_H / 2 - 0.10, 0.45, 0.20, "genes", size=8.5, color=C_BODY)
add_text(M_X, M_Y - 0.20, M_W, 0.15, "cells", size=8.5, color=C_BODY, align=PP_ALIGN.CENTER)

# label bar
LB_Y = M_Y + M_H + 0.15
add_text(A_X + 0.10, LB_Y - 0.18, COL_W - 0.20, 0.15, "cell labels", size=8.5, color=C_BODY)
rare_positions = {3, 11, 18}
for c in range(COLS):
    is_rare = c in rare_positions
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(M_X + c * cell_w), Inches(LB_Y),
        Inches(cell_w), Inches(0.18),
    )
    rect.shadow.inherit = False
    rect.fill.solid()
    rect.fill.fore_color.rgb = C_RARE if is_rare else C_MAJORITY
    rect.line.color.rgb = C_WHITE
    rect.line.width = Pt(0.5)

# legend
LG_Y = LB_Y + 0.32
add_rect(M_X, LG_Y, 0.13, 0.13, fill=C_MAJORITY, line=None)
add_text(M_X + 0.18, LG_Y - 0.02, 1.2, 0.18, "majority labelled cells", size=8.5, color=C_BODY)
add_rect(M_X + 1.6, LG_Y, 0.13, 0.13, fill=C_RARE, line=None)
add_text(M_X + 1.78, LG_Y - 0.02, 1.5, 0.18, "rare cells (~1–10%)", size=8.5, color=C_BODY)

# bottom note
add_text(A_X + 0.10, A_Y + PANEL_H - 0.32, COL_W - 0.20, 0.20,
         "Heavily imbalanced labels: only a few rare cells observed during training.",
         size=9, color=C_TITLE)

# ===============================================================
# PANEL B  Frozen scANVI backbone (latent scatter)
# ===============================================================
B_X, B_Y = COL1_X, PANEL_BOT_Y
panel_frame(B_X, B_Y, COL_W, PANEL_H, "B", "Frozen scANVI backbone",
            "semi-supervised encoder → latent z (10D; 2D UMAP for display)")

# plot area
PX, PY, PW, PH = B_X + 0.55, B_Y + 0.65, COL_W - 0.80, PANEL_H - 1.20
# axes
add_line(PX, PY + PH, PX + PW, PY + PH, color=C_AXIS, lw=0.75)  # x-axis
add_line(PX, PY, PX, PY + PH, color=C_AXIS, lw=0.75)            # y-axis
add_text(PX + PW - 0.40, PY + PH + 0.02, 0.50, 0.16, "UMAP 1", size=7, color=C_MUTED)
add_text(PX - 0.55, PY, 0.45, 0.16, "UMAP 2", size=7, color=C_MUTED)

# 3 majority clusters + 1 rare cluster
def scatter_cluster(cx, cy, n, spread, color, dot_d=0.05, seed=0):
    random.seed(seed)
    for _ in range(n):
        dx = random.gauss(0, spread)
        dy = random.gauss(0, spread)
        add_ellipse(cx + dx, cy + dy, dot_d, fill=color)


scatter_cluster(PX + 0.45, PY + 0.65, 16, 0.10, C_MAJORITY, seed=1)
scatter_cluster(PX + 1.55, PY + 1.45, 18, 0.11, C_MAJORITY, seed=2)
scatter_cluster(PX + 2.50, PY + 0.55, 14, 0.10, C_MAJORITY, seed=3)
scatter_cluster(PX + 1.20, PY + 0.20, 7, 0.07, C_RARE, dot_d=0.055, seed=4)

# cluster annotations
add_text(PX + 0.20, PY + 1.00, 0.9, 0.16, "majority A", size=8, color=C_MUTED)
add_text(PX + 1.40, PY + 1.78, 0.9, 0.16, "majority B", size=8, color=C_MUTED)
add_text(PX + 2.45, PY + 0.95, 0.9, 0.16, "majority C", size=8, color=C_MUTED)
add_text(PX + 1.08, PY + 0.38, 0.6, 0.16, "rare", size=8.5, bold=True, color=C_RARE)

# frozen badge
badge_x = B_X + COL_W - 1.10
add_rect(badge_x, B_Y + 0.40, 0.95, 0.22, transparent_fill=True, line=C_PHASE, line_w=0.6)
add_text(badge_x, B_Y + 0.40, 0.95, 0.22, "frozen backbone", size=8, color=C_TITLE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ===============================================================
# PANEL C  Prototype reference
# ===============================================================
C_PX, C_PY = COL2_X, PANEL_TOP_Y
panel_frame(C_PX, C_PY, COL_W, PANEL_H, "C", "Prototype reference",
            "class means + radii in latent z; separability score S")

PX, PY, PW, PH = C_PX + 0.55, C_PY + 0.65, COL_W - 0.80, PANEL_H - 1.20
add_line(PX, PY + PH, PX + PW, PY + PH, color=C_AXIS, lw=0.75)
add_line(PX, PY, PX, PY + PH, color=C_AXIS, lw=0.75)

# majority cluster 1
scatter_cluster(PX + 0.45, PY + 1.20, 12, 0.10, C_MAJORITY, seed=5)
add_ring(PX + 0.45, PY + 1.20, 0.55, line=C_MUTED, dash=True)
add_x_mark(PX + 0.45, PY + 1.20, size=0.10, color=C_LATENT, lw=2.2)

# majority cluster 2
scatter_cluster(PX + 2.50, PY + 1.55, 12, 0.10, C_MAJORITY, seed=6)
add_ring(PX + 2.50, PY + 1.55, 0.58, line=C_MUTED, dash=True)
add_x_mark(PX + 2.50, PY + 1.55, size=0.10, color=C_LATENT, lw=2.2)

# rare cluster
scatter_cluster(PX + 1.50, PY + 0.30, 6, 0.07, C_RARE, dot_d=0.055, seed=7)
add_ring(PX + 1.50, PY + 0.30, 0.35, line=C_RARE, dash=True)
add_x_mark(PX + 1.50, PY + 0.30, size=0.10, color=C_RARE, lw=2.2)

# separability line
add_line(PX + 0.45, PY + 1.20, PX + 1.50, PY + 0.30, color=C_BODY, lw=0.75, dash=True)
add_text(PX + 0.70, PY + 0.65, 1.2, 0.16, "d(rare, majority)", size=8, color=C_BODY)

# labels
add_text(PX + 0.85, PY + 1.30, 0.7, 0.16, "radius r", size=8, color=C_MUTED)
add_text(PX + 1.70, PY + 0.22, 1.0, 0.16, "rare prototype", size=8, bold=True, color=C_RARE)
add_text(PX + 0.20, PY + 1.62, 0.9, 0.16, "prototype", size=8, color=C_TITLE)

# separability formula
add_text(C_PX + 0.10, C_PY + PANEL_H - 0.32, COL_W - 0.20, 0.20,
         "S = d(rare, nearest majority) / mean intra-rare radius",
         size=9, color=C_TITLE)


# ===============================================================
# PANEL D  Candidate filtering (funnel)
# ===============================================================
D_X, D_Y = COL2_X, PANEL_BOT_Y
panel_frame(D_X, D_Y, COL_W, PANEL_H, "D", "Candidate filtering",
            "cells whose nearest prototype is the rare class become candidates")

# funnel top trapezoid
fx_left, fx_right = D_X + 0.30, D_X + COL_W - 1.20
fy_top = D_Y + 0.55
fy_mid = D_Y + 1.65
inner_left = D_X + 0.95
inner_right = D_X + COL_W - 1.85
funnel_top = [
    (fx_left, fy_top), (fx_right, fy_top),
    (inner_right, fy_mid), (inner_left, fy_mid),
]
add_freeform(funnel_top, fill=C_FILL_BG, fill_alpha=1.0, line=C_PHASE, line_w=0.75)
# output box (small trapezoid below)
out_y_top = fy_mid
out_y_bot = D_Y + 2.45
out_left_top, out_right_top = inner_left, inner_right
out_left_bot, out_right_bot = inner_left + 0.18, inner_right - 0.18
add_freeform(
    [(out_left_top, out_y_top), (out_right_top, out_y_top),
     (out_right_bot, out_y_bot), (out_left_bot, out_y_bot)],
    fill=C_FILL_BG, fill_alpha=1.0, line=C_PHASE, line_w=0.75,
)

# many gray dots inside top funnel
random.seed(11)
for _ in range(36):
    rx = random.uniform(fx_left + 0.05, fx_right - 0.05)
    ry = random.uniform(fy_top + 0.05, fy_mid - 0.10)
    add_ellipse(rx, ry, 0.045, fill=C_MAJORITY)
# a few rare dots scattered
for cx, cy in [(D_X + 1.10, D_Y + 0.85), (D_X + 1.70, D_Y + 1.10), (D_X + 1.40, D_Y + 1.40)]:
    add_ellipse(cx, cy, 0.06, fill=C_RARE)

# bottom output dots — mostly rare + 1-2 gray
random.seed(12)
for cx, cy in [
    (inner_left + 0.30, out_y_top + 0.40),
    (inner_left + 0.70, out_y_top + 0.55),
    (inner_left + 1.10, out_y_top + 0.40),
    (inner_left + 1.50, out_y_top + 0.55),
]:
    add_ellipse(cx, cy, 0.065, fill=C_RARE)
add_ellipse(inner_left + 0.50, out_y_top + 0.55, 0.05, fill=C_MAJORITY)
add_ellipse(inner_left + 1.30, out_y_top + 0.40, 0.05, fill=C_MAJORITY)

# side labels
side_x = D_X + COL_W - 1.10
for i, label in enumerate(["nearest prototype", "rank / margin", "safety check"]):
    yy = D_Y + 0.75 + i * 0.45
    add_line(side_x - 0.10, yy, side_x, yy, color=C_BODY, lw=0.8)
    add_text(side_x + 0.02, yy - 0.10, 1.1, 0.20, label, size=8, color=C_BODY)

# top/bottom captions
add_text(D_X + 0.30, fy_top - 0.20, COL_W - 0.6, 0.15,
         "all unlabelled / mis-predicted cells", size=8, color=C_MUTED)
add_text(D_X + 0.30, D_Y + PANEL_H - 0.30, COL_W - 0.6, 0.20,
         "candidate cells", size=9, color=C_TITLE, align=PP_ALIGN.CENTER)


# ===============================================================
# PANEL E  Membership score + conformal τ
# ===============================================================
E_X, E_Y = COL3_X, PANEL_TOP_Y
panel_frame(E_X, E_Y, COL_W, PANEL_H, "E", "Membership score and conformal τ",
            "validation distribution of rare-membership scores, calibrated to FFR ≤ α")

PX, PY, PW, PH = E_X + 0.55, E_Y + 0.65, COL_W - 0.80, PANEL_H - 1.20
add_line(PX, PY + PH, PX + PW, PY + PH, color=C_AXIS, lw=0.75)
add_line(PX, PY, PX, PY + PH, color=C_AXIS, lw=0.75)
add_text(PX, PY + PH + 0.05, PW, 0.18, "rare membership score", size=9, color=C_TITLE, align=PP_ALIGN.CENTER)
add_text(PX - 0.55, PY + PH / 2 - 0.10, 0.45, 0.16, "density", size=8, color=C_MUTED)


# distribution shapes — non-rare bell (gray, peak left)
def bell_points(x0, x1, peak_x, peak_y, base_y, n=24):
    """Return polygon approx of a Gaussian bump from x0..x1, peak at peak_x."""
    sigma = (x1 - x0) / 5.0
    pts = [(x0, base_y)]
    for i in range(n + 1):
        xi = x0 + i * (x1 - x0) / n
        h = (peak_y - base_y) * math.exp(-((xi - peak_x) ** 2) / (2 * sigma ** 2))
        pts.append((xi, base_y - h))
    pts.append((x1, base_y))
    return pts


base = PY + PH
non_rare = bell_points(PX, PX + PW * 0.75, PX + PW * 0.30, base - 1.55, base, n=30)
rare_pts = bell_points(PX + PW * 0.55, PX + PW * 1.00, PX + PW * 0.83, base - 1.20, base, n=24)
add_freeform(non_rare, fill=C_MAJORITY, fill_alpha=0.6, line=C_PHASE, line_w=0.75)
add_freeform(rare_pts, fill=C_RARE, fill_alpha=0.55, line=C_RARE, line_w=0.75)

# threshold τ
tau_x = PX + PW * 0.72
add_line(tau_x, PY, tau_x, base, color=C_TAU, lw=1.6, dash=True)
add_text(tau_x + 0.05, PY, 0.6, 0.20, "τ", size=14, bold=True, color=C_TAU)
add_text(tau_x + 0.05, PY + 0.22, 1.2, 0.18, "calibrated threshold", size=8, color=C_TAU)
add_text(tau_x + 0.05, PY + 0.40, 1.2, 0.18, "α = 0.01", size=8, color=C_TAU)

# legend
add_rect(PX, PY - 0.20, 0.13, 0.12, fill=C_MAJORITY, line=None)
add_text(PX + 0.18, PY - 0.22, 1.6, 0.18, "non-rare val cells", size=8, color=C_BODY)
add_rect(PX + 1.50, PY - 0.20, 0.13, 0.12, fill=C_RARE, line=None)
add_text(PX + 1.68, PY - 0.22, 1.4, 0.18, "rare val cells", size=8, color=C_BODY)


# ===============================================================
# PANEL F  Refined prediction (before / after)
# ===============================================================
F_X, F_Y = COL3_X, PANEL_BOT_Y
panel_frame(F_X, F_Y, COL_W, PANEL_H, "F", "Refined prediction",
            "candidates with score ≥ τ are relabelled to the rare class")

mini_w = (COL_W - 0.80) / 2
mini_h = PANEL_H - 1.10
miniA_x = F_X + 0.25
miniB_x = F_X + COL_W - 0.25 - mini_w
mini_y = F_Y + 0.65

for title, mx in [("before (baseline)", miniA_x), ("after (rescued)", miniB_x)]:
    add_text(mx, mini_y - 0.22, mini_w, 0.18, title, size=9, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
    add_rect(mx, mini_y, mini_w, mini_h, fill=C_FILL_BG, line=C_FRAME, line_w=0.5)
    # axes
    add_line(mx + 0.10, mini_y + mini_h - 0.05, mx + mini_w - 0.05, mini_y + mini_h - 0.05, color=C_AXIS, lw=0.5)
    add_line(mx + 0.10, mini_y + 0.05, mx + 0.10, mini_y + mini_h - 0.05, color=C_AXIS, lw=0.5)

# BEFORE — majority cluster + a few orange (visible at edge) + rare missed shown grey at top
scatter_cluster(miniA_x + mini_w * 0.30, mini_y + mini_h * 0.55, 12, 0.07, C_MAJORITY, seed=21)
scatter_cluster(miniA_x + mini_w * 0.65, mini_y + mini_h * 0.80, 10, 0.06, C_MAJORITY, seed=22)
# 2 visible rare (orange, near rest of rare cluster region)
add_ellipse(miniA_x + mini_w * 0.55, mini_y + mini_h * 0.20, 0.06, fill=C_RARE)
add_ellipse(miniA_x + mini_w * 0.65, mini_y + mini_h * 0.25, 0.06, fill=C_RARE)
# 3 missed rare drawn grey
for dx, dy in [(0.40, 0.12), (0.48, 0.18), (0.36, 0.20)]:
    add_ellipse(miniA_x + mini_w * dx, mini_y + mini_h * dy, 0.06, fill=C_MAJORITY)
add_text(miniA_x + 0.10, mini_y + mini_h - 0.30, mini_w - 0.20, 0.18,
         "3 rare missed", size=7.5, color=C_MUTED)

# arrow between
arrow_y = mini_y + mini_h / 2
add_line(miniA_x + mini_w + 0.06, arrow_y, miniB_x - 0.06, arrow_y, color=C_BODY, lw=1.4, arrow=True)
add_text(miniA_x + mini_w + 0.05, arrow_y - 0.22, miniB_x - miniA_x - mini_w - 0.10, 0.16,
         "conformal", size=8, color=C_BODY, align=PP_ALIGN.CENTER)
add_text(miniA_x + mini_w + 0.05, arrow_y + 0.08, miniB_x - miniA_x - mini_w - 0.10, 0.16,
         "rescue", size=8, color=C_BODY, align=PP_ALIGN.CENTER)

# AFTER — same layout, missed rare now green (rescued)
scatter_cluster(miniB_x + mini_w * 0.30, mini_y + mini_h * 0.55, 12, 0.07, C_MAJORITY, seed=23)
scatter_cluster(miniB_x + mini_w * 0.65, mini_y + mini_h * 0.80, 10, 0.06, C_MAJORITY, seed=24)
add_ellipse(miniB_x + mini_w * 0.55, mini_y + mini_h * 0.20, 0.06, fill=C_RARE)
add_ellipse(miniB_x + mini_w * 0.65, mini_y + mini_h * 0.25, 0.06, fill=C_RARE)
for dx, dy in [(0.40, 0.12), (0.48, 0.18), (0.36, 0.20)]:
    add_ellipse(miniB_x + mini_w * dx, mini_y + mini_h * dy, 0.07, fill=C_RESCUED)
add_text(miniB_x + 0.10, mini_y + mini_h - 0.30, mini_w - 0.20, 0.18,
         "3 rescued, FFR ≤ α", size=7.5, color=C_MUTED)

# legend at bottom
leg_y = F_Y + PANEL_H - 0.36
add_rect(F_X + 0.20, leg_y, 0.13, 0.13, fill=C_MAJORITY, line=None)
add_text(F_X + 0.38, leg_y - 0.02, 0.8, 0.18, "majority", size=8, color=C_BODY)
add_rect(F_X + 1.15, leg_y, 0.13, 0.13, fill=C_RARE, line=None)
add_text(F_X + 1.33, leg_y - 0.02, 0.8, 0.18, "true rare", size=8, color=C_BODY)
add_rect(F_X + 2.10, leg_y, 0.13, 0.13, fill=C_RESCUED, line=None)
add_text(F_X + 2.28, leg_y - 0.02, 1.2, 0.18, "rescued rare", size=8, color=C_BODY)


# ===============================================================
# Inter-phase arrows (horizontal) and intra-column arrows (vertical)
# ===============================================================
mid_y = (PANEL_TOP_Y + PANEL_H + PANEL_BOT_Y) / 2

# Phase 1 -> Phase 2 horizontal arrow
add_line(COL1_X + COL_W + 0.02, mid_y, COL2_X - 0.05, mid_y, color=C_BODY, lw=1.4, arrow=True)
# Phase 2 -> Phase 3
add_line(COL2_X + COL_W + 0.02, mid_y, COL3_X - 0.05, mid_y, color=C_BODY, lw=1.4, arrow=True)

# A -> B vertical arrow
mid_x_col1 = COL1_X + COL_W / 2
add_line(mid_x_col1, PANEL_TOP_Y + PANEL_H + 0.02, mid_x_col1, PANEL_BOT_Y - 0.04, color=C_ARROW, lw=1.0, arrow=True)
# C -> D
mid_x_col2 = COL2_X + COL_W / 2
add_line(mid_x_col2, PANEL_TOP_Y + PANEL_H + 0.02, mid_x_col2, PANEL_BOT_Y - 0.04, color=C_ARROW, lw=1.0, arrow=True)
# E -> F
mid_x_col3 = COL3_X + COL_W / 2
add_line(mid_x_col3, PANEL_TOP_Y + PANEL_H + 0.02, mid_x_col3, PANEL_BOT_Y - 0.04, color=C_ARROW, lw=1.0, arrow=True)

# bottom note
add_text(M, SLIDE_H - 0.30, SLIDE_W - 2 * M, 0.18,
         "All gates and thresholds are calibrated on train / validation only; test labels are used solely for final evaluation.",
         size=8.5, color=C_MUTED)


prs.save(str(OUT))
print(f"[saved] {OUT}")
